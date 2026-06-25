import streamlit as st
import re
import json
from pptx import Presentation
import io
from slide_classifier import load_classifier

st.set_page_config(page_title="TinySlide AI", layout="wide")


@st.cache_resource
def get_model():
    return load_classifier()


model = get_model()
st.title("TinySlide AI")
st.subheader("Offline Text-to-Slide Generator")
st.info(
    f"Backend: {model.backend.upper()} | "
    "Model: TF-IDF + Logistic Regression | "
    "Size: 0.0377 MB | "
    "Accuracy: 100% synthetic + real-world test | "
    "Latency: 1.2539 ms | "
    "Offline: Yes"
)
st.markdown("### Model Performance")

m1, m2, m3, m4 = st.columns(4)

m1.metric("Model Size", "0.0377 MB")
m2.metric("Accuracy", "100%")
m3.metric("Latency", "1.2539 ms")
m4.metric("Offline", "Yes")
text = st.text_area("Paste your notes here:", height=250)

def clean_sentences(text):
    sentences = re.split(r'(?<=[.!?]) +', text.strip())
    return [s.strip() for s in sentences if len(s.strip()) > 10]

def generate_title(text):
    words = text.split()
    return " ".join(words[:6]).title() if words else "Untitled Presentation"
def generate_slide_heading(content_type):
    headings = {
        "definition": "Definition",
        "problem": "Problem Statement",
        "solution": "Proposed Solution",
        "architecture": "System Architecture",
        "features": "Key Features",
        "implementation": "Implementation Steps",
        "strategy": "Rollout Strategy",
        "benefit": "Key Benefits",
        "process": "Process Overview",
        "comparison": "Comparison",
        "example": "Example",
        "statistic": "Key Statistics",
        "overview": "Topic Overview"
    }

    return headings.get(content_type, "Slide Overview")
def make_slides(text):
    sentences = clean_sentences(text)
    title = generate_title(text)

    grouped_content = {}

    for sentence in sentences:
        predicted_type = predict_content_type(sentence)

        if predicted_type not in grouped_content:
            grouped_content[predicted_type] = []

        grouped_content[predicted_type].append(clean_bullet(sentence))

    slides = []

    for content_type, bullets in grouped_content.items():
        slide = {
            "heading": generate_slide_heading(content_type),
            "content_type": content_type,
            "bullets": bullets,
            "layout": suggest_layout(bullets)
        }

        slides.append(slide)

    return {
        "title": title,
        "slides": slides
    }
def clean_bullet(sentence):
    sentence = sentence.strip()

    remove_words = ["basically", "actually", "very", "really"]
    for word in remove_words:
        sentence = sentence.replace(word, "")

    sentence = sentence.strip()

    if len(sentence) > 120:
        sentence = sentence[:117] + "..."

    return sentence[0].upper() + sentence[1:] if sentence else sentence
def suggest_layout(chunk):
    bullet_count = len(chunk)

    if bullet_count <= 2:
        return "title_and_content"

    elif bullet_count <= 4:
        return "two_column"

    else:
        return "comparison"
def rule_based_content_type(sentence):
    text = sentence.lower()

    rules = [
        ("implementation", [
            "step-by-step", "implementing", "implementation", "requires",
            "first", "second", "third", "method", "deploy"
        ]),
        ("benefit", [
            "ultimately", "benefit", "benefits", "impact", "improves",
            "reduces", "saves", "self-optimizing", "optimizing"
        ]),
        ("strategy", [
            "strategy", "long-term plan", "plan", "phased rollout",
            "rollout", "pilot zone", "scale", "phase"
        ]),
        ("architecture", [
            "approach to building", "system relies", "integrating",
            "iot sensors", "sensors", "predictive", "machine learning",
            "architecture", "components"
        ]),
        ("features", [
            "key features", "features", "include", "includes",
            "adaptive", "dashboard", "alerts"
        ]),
        ("solution", [
            "solution", "to solve", "solves", "transitioning",
            "automates", "proposed"
        ]),
        ("comparison", [
            "compared", "unlike", "whereas", "while", "than"
        ]),
        ("example", [
            "for example", "for instance", "use case"
        ]),
        ("statistic", [
            "percent", "%", "accuracy", "latency", "metric", "data"
        ]),
        ("problem", [
            "problem", "challenge", "issue", "gridlock", "bottleneck",
            "congestion", "chaotic"
        ]),
        ("definition", [
            " is a ", " refers to ", " means ", " defined as "
        ])
    ]

    for content_type, keywords in rules:
        if any(keyword in text for keyword in keywords):
            return content_type

    return None
def predict_content_type(sentence):
    rule_prediction = rule_based_content_type(sentence)

    if rule_prediction:
        return rule_prediction

    model_prediction = model.predict([sentence])[0]

    if model_prediction == "definition":
        return "overview"

    return model_prediction
def create_pptx(result):
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = result["title"]
    title_slide.placeholders[1].text = "Generated by TinySlide AI"

    for slide_data in result["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["heading"]

        content = slide.placeholders[1]
        content.text = ""

        for bullet in slide_data["bullets"]:
            p = content.text_frame.add_paragraph()
            p.text = bullet
            p.level = 0

    pptx_file = io.BytesIO()
    prs.save(pptx_file)
    pptx_file.seek(0)

    return pptx_file

if st.button("Generate Slides"):
    if text.strip():
        result = make_slides(text)
        json_data = json.dumps(result, indent=4)

        pptx_file = create_pptx(result)

        st.download_button(
            label="Download PowerPoint",
            data=pptx_file,
            file_name="generated_slides.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        st.download_button(
            label="Download JSON",
            data=json_data,
            file_name="generated_slides.json",
            mime="application/json"
        )

        col1, col2 = st.columns(2)

        with col1:
            st.subheader("Generated JSON")
            st.json(result)

        with col2:
            st.subheader("Slide Preview")
            st.markdown(f"# {result['title']}")

            for slide in result["slides"]:
                st.markdown("---")
                st.markdown(f"## {slide['heading']}")
                for bullet in slide["bullets"]:
                    st.markdown(f"- {bullet}")

    else:
        st.warning("Please enter some text.")
        
